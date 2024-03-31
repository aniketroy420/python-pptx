from pptx import Presentation
from pptx.util import Inches, Pt,Cm
from pptx.dml.color import RGBColor
#from pptx.enum.shapes import MSO_SHAPE

from pptx.enum.text import PP_ALIGN

prs = Presentation()

prs.slide_width = Cm(33.9)
prs.slide_height = Cm(19.07)


blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)


# this sets backgrund color of the slide
# background = slide.background
# fill = background.fill
# fill.solid()
# fill.fore_color.rgb = RGBColor(37, 150, 190)

#  image path
logo_path = "logo.png"
paytm_img = "paytm.png"  
lookbeyond_img = "lookbeyond.png"
logo_left = 0.5  # Adjust these values to reposition the image
top_img = 0.5
paytm_left = 27.34


#This sets the background image of the slide
left = top = 0
img_path = "bg.jpg"
pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)


# Adds the image 
logo = slide.shapes.add_picture(logo_path, left=Cm(1.57),top= Cm(1.3))
paytm = slide.shapes.add_picture(paytm_img,left=Cm(27.34), top= Cm(1.67))
lookbeyond= slide.shapes.add_picture(lookbeyond_img,left=0,top = Cm(16.53))

## Sets textbox position 
left = Cm(1.43)
top = Cm(8.05)
width = Cm(13)
height = Cm(2.94)

# Creates textbox object
textbox = slide.shapes.add_textbox(left, top,width,height)

# Gets text frame (to access text properties)
tf = textbox.text_frame
tf.word_wrap = True
# Add your text here
text_content = [
    "PAYTM JUNE 2023"
]
textbox.text_frame.text = "CONFIDENTIAL"
textbox.text_frame.word_wrap = True
p = textbox.text_frame.paragraphs[0]  # Access first paragraph
p.font.name = 'Calibri'  # Set font to Calibri
p.font.bold = True  # Set bold formatting
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.size = Pt(43.5)
p.alignment = PP_ALIGN.LEFT
#Use different properties for the different text available in the above list
for i in text_content:
    p = tf.add_paragraph()
    p.font.bold = True
    p.text = i
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(17.5)

# 2nd slide----------------------------------------------------------------------------------------
        
slide = prs.slides.add_slide(prs.slide_layouts[6])

logo_blue_path = "logo_blue.png"
logo_blue = slide.shapes.add_picture(logo_blue_path, left = Cm(1.46),top= Cm(4.52))

r_semi = "r_semi.png"
r_semi = slide.shapes.add_picture(r_semi, left = Cm(2.35),top= Cm(15.15))

# Creates textbox object
textbox = slide.shapes.add_textbox(left=Cm(1.46), top= Cm(3.25),width=Cm(3.49),height= Cm(1.13))
textbox.text_frame.text = "We are"
textbox.text_frame.word_wrap = True
p = textbox.text_frame.paragraphs[0]  # Access first paragraph
p.font.name = 'Calibri'  # Set font to Calibri
p.font.bold = True  # Set bold formatting
p.font.color.rgb = RGBColor(102, 46, 137)
p.font.size = Pt(25)
p.alignment = PP_ALIGN.LEFT


# fIRST paragraph ----------------------------------------------------------------------------

textbox = slide.shapes.add_textbox(left= Cm(1.43), top= Cm(7.16),width= Cm(15.4),height=Cm(5.37))
tf = textbox.text_frame
tf.word_wrap = True

textbox.text_frame.text = "Executive Access India is a well-established executive search firm that has been operating for over 28 years in the country. We have contributed significantly to our clients' success by building world-class leadership teams across industry verticals.\n"
textbox.text_frame.word_wrap = True
p = textbox.text_frame.paragraphs[0]  # Access first paragraph
p.font.name = 'Calibri'  # Set font to Calibri
p.font.color.rgb = RGBColor(89,89,91)
p.font.size = Pt(12)
p.alignment = PP_ALIGN.JUSTIFY

first_para = [
    "We are a member of the Panorama community of over 200 professionals and 22 partners spread across North America, South America, Europe, and Asia Pacific. Our clients include Fortune 500 companies, mid-sized global companies, as well as some of the most respected Indian corporates and start-ups."
]
for i in first_para:
    p = tf.add_paragraph()
    p.text = i
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(89,89,91)
    p.alignment = PP_ALIGN.JUSTIFY
    p.font.size = Pt(12)

# second para --------------------------------------------------------------------------------------------------
textbox = slide.shapes.add_textbox(left= Cm(17.35), top= Cm(7.17),width= Cm(15.13),height=Cm(7.89))
tf = textbox.text_frame
tf.word_wrap = True

textbox.text_frame.text = "One of our core strengths is our ability to simultaneously be both global and local. As a boutique firm, we are dedicated to giving our clients and candidates undivided attention, ensuring that every search is tailored to their unique needs and requirements.\n"
textbox.text_frame.word_wrap = True
p = textbox.text_frame.paragraphs[0]  # Access first paragraph
p.font.name = 'Calibri'  # Set font to Calibri
p.font.color.rgb = RGBColor(89,89,91)
p.font.size = Pt(12)
p.alignment = PP_ALIGN.JUSTIFY

second_para = [
    "The company's logo is derived from the symbol \"Ren\", which signifies the plenitude of humanness. Our commitment to our clients is reflected in our values of Client First, Integrity, Collaboration, and Nurturing.",
    "",
    "We are the pioneers of the “Accountability Clause” in the Indian retained search market and have always striven to align our business model with the needs of our clients."
]
for i in second_para:
    p = tf.add_paragraph()
    p.text = i
    p.font.name = 'Calibri'
    p.font.color.rgb = RGBColor(89,89,91)
    p.alignment = PP_ALIGN.JUSTIFY
    p.font.size = Pt(12)

# 3rd para on the semi circle

textbox = slide.shapes.add_textbox(left=Cm(3.47), top= Cm(15.31),width=Cm(5.62),height= Cm(2.6))
textbox.text_frame.text = "“We have learnt that by uniting the right people with the right corporations, we can help both to achieve their full potential.”"
textbox.text_frame.word_wrap = True
p = textbox.text_frame.paragraphs[0]  # Access first paragraph
p.font.name = 'Calibri'  # Set font to Calibri
p.font.color.rgb = RGBColor(255,255,255)
p.font.size = Pt(11)
p.alignment = PP_ALIGN.CENTER

# 3nd slide----------------------------------------------------------------------------------------
slide = prs.slides.add_slide(prs.slide_layouts[6])

#images
left = top = 0
img_path = "bg2.png"
bg = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)
bg_asset = "bg_asset.png"
asset = slide.shapes.add_picture(bg_asset, left = Cm(12.31),top= Cm(10.09))
content = "content.png"
content_img = slide.shapes.add_picture(content, left = 0,top= Cm(15.77))

#text
x = 1.12
y = 1.12
text = 1
for i in range(7):
    if(i<4):
        textbox = slide.shapes.add_textbox(left=Cm(1.65), top= Cm(x),width=Cm(3.56),height= Cm(2.82))
        textbox.text_frame.text = str(text)
        textbox.text_frame.word_wrap = True
        p = textbox.text_frame.paragraphs[0]  # Access first paragraph
        p.font.name = 'Arial'  # Set font to Calibri
        p.font.bold = True  # Set bold formatting
        p.font.color.rgb = RGBColor(165, 165, 165)
        p.font.size = Pt(60)
        p.alignment = PP_ALIGN.LEFT
        x +=2.6
        
    else:
        textbox = slide.shapes.add_textbox(left=Cm(13.64), top= Cm(y),width=Cm(3.56),height= Cm(2.82))
        textbox.text_frame.text = str(text)
        textbox.text_frame.word_wrap = True
        p = textbox.text_frame.paragraphs[0]  # Access first paragraph
        p.font.name = 'Arial'  # Set font to Calibri
        p.font.bold = True  # Set bold formatting
        p.font.color.rgb = RGBColor(165, 165, 165)
        p.font.size = Pt(60)
        p.alignment = PP_ALIGN.LEFT
        y +=2.6
    text += 1
text_content = ["Our Understanding of the Role",
                "Ideal Candidate Profile",
                "Assessment Criteria",
                "Target Talent Landscape",
                "Indicative Target Universe",
                "Benchmark Profiles",
                "Search Process & Timelines"

]
top1 = 2.61
top2 = 2.61
for i in range(len(text_content)):
            if(i<4):
                textbox = slide.shapes.add_textbox(left=Cm(2.27), top= Cm(top1),width=Cm(9.8),height= Cm(1.03))
                textbox.text_frame.text = text_content[i]
                textbox.text_frame.word_wrap = True
                p = textbox.text_frame.paragraphs[0]  # Access first paragraph
                p.font.name = 'Arial'  # Set font to Calibri
                p.font.bold = True
                p.font.underline = True  # Set bold formatting
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
                top1 +=2.6
            else:
                textbox = slide.shapes.add_textbox(left=Cm(14.26), top= Cm(top2),width=Cm(9.23),height= Cm(1.03))
                textbox.text_frame.text = text_content[i]
                textbox.text_frame.word_wrap = True
                p = textbox.text_frame.paragraphs[0]  # Access first paragraph
                p.font.name = 'Arial'  # Set font to Calibri
                p.font.bold = True  # Set bold formatting
                p.font.underline = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
                top2 +=2.6
                 

print(prs.slide_width)
print(prs.slide_height)
# Save the presentation
prs.save("4.pptx")
